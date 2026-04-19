"""
agents/delivery.py

Delivery node: format SynthesisOutput into the requested output format.

Reads:  state["synthesis"], state["critic_output"], state["output_format"],
        state["revision_count"]
Writes: state["final_output"]

Binary formats (DOCX, PDF, PPTX) are rendered inside an isolated microsandbox
container and returned as base64-encoded data URIs so the existing
Optional[str] field in ResearchAgentState requires no schema change.
"""

from __future__ import annotations

import base64
import json

from agents.sandbox_runner import execute_in_sandbox
from graph.progress import delivery_done, delivery_start
from graph.state import (
    CriticOutput,
    Evidence,
    Finding,
    ResearchAgentState,
    SynthesisOutput,
)

_MIME_TYPES: dict[str, str] = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_QUALITY_CAVEAT_TEMPLATE: str = (
    "> **QUALITY NOTICE:** This research brief was delivered after reaching "
    "the maximum revision limit (2 revisions) without meeting the quality "
    "threshold (score: {score:.2f}/5.0). Review with appropriate scrutiny.\n"
)

# Must match graph/router.py._MAX_REVISIONS — declared locally to avoid
# circular import risk (router imports from state, state imports nothing from agents).
_MAX_REVISIONS: int = 2


# ---------------------------------------------------------------------------
# Force-delivery detection
# ---------------------------------------------------------------------------


def _is_force_delivered(state: ResearchAgentState) -> bool:
    """Return True if revision_count >= _MAX_REVISIONS and critic did not pass.

    Args:
        state: Current ResearchAgentState.

    Returns:
        True when max revisions exhausted without meeting quality threshold.
    """
    critic_output: CriticOutput | None = state["critic_output"]
    if critic_output is None:
        return False
    return state["revision_count"] >= _MAX_REVISIONS and not critic_output.passed


# ---------------------------------------------------------------------------
# Citation map
# ---------------------------------------------------------------------------


def _build_citation_map(citations: list[Evidence]) -> dict[str, int]:
    """Build URL → 1-based footnote index mapping from synthesis.citations.

    Args:
        citations: List of Evidence objects from SynthesisOutput.citations.

    Returns:
        Dict mapping source_url to 1-based footnote index.
    """
    return {ev.source_url: i for i, ev in enumerate(citations, start=1)}


# ---------------------------------------------------------------------------
# Finding renderer
# ---------------------------------------------------------------------------


def _render_finding(finding: Finding, citation_map: dict[str, int]) -> str:
    """Render a single finding as markdown.

    Includes: ### heading, body paragraph, [^N] footnote refs for each
    evidence_ref URL present in citation_map. Unknown URLs are silently
    skipped (no footnote appended).

    Args:
        finding: Finding dataclass to render.
        citation_map: URL → footnote index mapping.

    Returns:
        Markdown string for the finding block.
    """
    refs = "".join(
        f"[^{citation_map[url]}]"
        for url in finding.evidence_refs
        if url in citation_map
    )
    return f"### {finding.heading}\n\n{finding.body}{refs}"


# ---------------------------------------------------------------------------
# Citations section renderer
# ---------------------------------------------------------------------------


def _render_citations_section(citations: list[Evidence]) -> str:
    """Render the References section with numbered footnote definitions.

    Format per entry:
        [^1]: [Source Title](URL) - Published: YYYY-MM-DD

    Args:
        citations: List of Evidence objects from SynthesisOutput.citations.

    Returns:
        Full References section string, or empty string if citations is empty.
    """
    if not citations:
        return ""
    lines: list[str] = ["## References", ""]
    for i, ev in enumerate(citations, start=1):
        pub = f" - Published: {ev.published_date}" if ev.published_date else ""
        lines.append(f"[^{i}]: [{ev.source_title}]({ev.source_url}){pub}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Quality badge
# ---------------------------------------------------------------------------


def _render_quality_badge(
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Render the quality badge line.

    Args:
        critic_output: CriticOutput from state, or None if critic never ran.
        is_force_delivered: True when max revisions reached without passing.

    Returns:
        Badge string, e.g. "**Quality: PASS** (4.25/5.0)"
    """
    if critic_output is None:
        return "**Quality: N/A**"
    score = critic_output.final_quality_score
    if is_force_delivered:
        return f"**Quality: REVIEW RECOMMENDED** ({score:.2f}/5.0)"
    return f"**Quality: PASS** ({score:.2f}/5.0)"


# ---------------------------------------------------------------------------
# Markdown renderer (exported for testing and reuse)
# ---------------------------------------------------------------------------


def render_markdown(
    synthesis: SynthesisOutput,
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Render SynthesisOutput as a complete markdown document.

    Pure function: no side effects, no LLM calls, deterministic output.

    Section order:
        1. Quality caveat (ONLY if force-delivered)
        2. # {topic}
        3. Quality badge line
        4. ## Executive Summary
        5. ## Key Findings (### per finding, [^N] footnote refs)
        6. ## Risks (omit if empty)
        7. ## Known Gaps (omit if empty)
        8. ## References (footnote definitions, omit if no citations)

    Args:
        synthesis: SynthesisOutput from the coordinator node.
        critic_output: CriticOutput from the critic node, or None.
        is_force_delivered: True when max revisions exhausted without passing.

    Returns:
        Complete markdown document as a single string.
    """
    citation_map = _build_citation_map(synthesis.citations)
    parts: list[str] = []

    # 1. Quality caveat (force-delivery only)
    if is_force_delivered and critic_output is not None:
        caveat = _QUALITY_CAVEAT_TEMPLATE.format(
            score=critic_output.final_quality_score
        )
        parts.append(caveat)

    # 2. Title
    parts.append(f"# {synthesis.topic}")

    # 3. Quality badge
    parts.append(_render_quality_badge(critic_output, is_force_delivered))

    # 4. Executive summary
    parts.append(f"## Executive Summary\n\n{synthesis.executive_summary}")

    # 5. Key findings
    if synthesis.findings:
        finding_blocks = "\n\n".join(
            _render_finding(f, citation_map) for f in synthesis.findings
        )
        parts.append(f"## Key Findings\n\n{finding_blocks}")
    else:
        parts.append("## Key Findings")

    # 6. Risks (omit if empty)
    if synthesis.risks:
        risk_lines = "\n".join(f"- {r}" for r in synthesis.risks)
        parts.append(f"## Risks\n\n{risk_lines}")

    # 7. Known Gaps (omit if empty)
    if synthesis.gaps:
        gap_lines = "\n".join(f"- {g}" for g in synthesis.gaps)
        parts.append(f"## Known Gaps\n\n{gap_lines}")

    # 8. References (omit if no citations)
    citations_section = _render_citations_section(synthesis.citations)
    if citations_section:
        parts.append(citations_section)

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Binary output helper
# ---------------------------------------------------------------------------


def _encode_binary_output(data: bytes, format_key: str) -> str:
    """Encode binary document bytes as a data URI string.

    Args:
        data: Raw document bytes from sandbox.
        format_key: One of "docx", "pdf", "pptx" — must be a key in _MIME_TYPES.

    Returns:
        Data URI: "data:{mime};base64,{encoded}"
    """
    mime = _MIME_TYPES[format_key]
    encoded = base64.b64encode(data).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def _build_doc_data(
    synthesis: SynthesisOutput,
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Serialize synthesis to a JSON string for embedding in sandbox code.

    repr(json.dumps(...)) embeds synthesis data as a Python string literal.
    json.dumps normalises the data; repr() wraps it in quotes and escapes
    any characters that could break the outer f-string. Synthesis content
    originates from the coordinator LLM — not raw user input — but repr()
    ensures no synthesis text can escape the string boundary regardless.

    Args:
        synthesis: SynthesisOutput to serialise.
        critic_output: CriticOutput or None.
        is_force_delivered: True when max revisions exhausted.

    Returns:
        repr() of the JSON string, ready to embed as a Python literal.
    """
    payload = {
        "topic": synthesis.topic,
        "executive_summary": synthesis.executive_summary,
        "findings": [
            {"heading": f.heading, "body": f.body} for f in synthesis.findings
        ],
        "risks": synthesis.risks,
        "gaps": synthesis.gaps,
        "citations": [
            {"title": c.source_title, "url": c.source_url, "date": c.published_date}
            for c in synthesis.citations
        ],
        "is_force_delivered": is_force_delivered,
        "quality_score": (
            critic_output.final_quality_score if critic_output is not None else None
        ),
    }
    return repr(json.dumps(payload))


# ---------------------------------------------------------------------------
# Binary format renderers
# ---------------------------------------------------------------------------


async def _render_docx(
    synthesis: SynthesisOutput,
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Render SynthesisOutput as DOCX via sandboxed python-docx execution.

    Args:
        synthesis: SynthesisOutput from the coordinator node.
        critic_output: CriticOutput from the critic node, or None.
        is_force_delivered: True when max revisions exhausted without passing.

    Returns:
        Data URI string containing base64-encoded DOCX bytes.

    Raises:
        SandboxExecutionError: If sandbox execution or package install fails.
    """
    doc_data_repr = _build_doc_data(synthesis, critic_output, is_force_delivered)
    code = f"""
import json
from docx import Document
from docx.shared import Pt

data = json.loads({doc_data_repr})

doc = Document()

if data["is_force_delivered"] and data["quality_score"] is not None:
    p = doc.add_paragraph()
    run = p.add_run(
        f"QUALITY NOTICE: Score {{data['quality_score']:.2f}}/5.0 — "
        "review with appropriate scrutiny."
    )
    run.bold = True

doc.add_heading(data["topic"], level=0)
doc.add_heading("Executive Summary", level=1)
doc.add_paragraph(data["executive_summary"])

if data["findings"]:
    doc.add_heading("Key Findings", level=1)
    for f in data["findings"]:
        doc.add_heading(f["heading"], level=2)
        doc.add_paragraph(f["body"])

if data["risks"]:
    doc.add_heading("Risks", level=1)
    for r in data["risks"]:
        doc.add_paragraph(r, style="List Bullet")

if data["gaps"]:
    doc.add_heading("Known Gaps", level=1)
    for g in data["gaps"]:
        doc.add_paragraph(g, style="List Bullet")

if data["citations"]:
    doc.add_heading("References", level=1)
    for i, c in enumerate(data["citations"], 1):
        date_part = (" (" + c["date"] + ")") if c["date"] else ""
        doc.add_paragraph("[" + str(i) + "] " + c["title"] + date_part + " \u2014 " + c["url"])

doc.save("/tmp/out")
"""
    docx_bytes = await execute_in_sandbox(code=code, packages=["python-docx", "lxml"])
    return _encode_binary_output(docx_bytes, "docx")


async def _render_pdf(
    synthesis: SynthesisOutput,
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Render SynthesisOutput as PDF via WeasyPrint in a sandbox.

    Converts the markdown render to HTML then passes it to WeasyPrint.

    Args:
        synthesis: SynthesisOutput from the coordinator node.
        critic_output: CriticOutput from the critic node, or None.
        is_force_delivered: True when max revisions exhausted without passing.

    Returns:
        Data URI string containing base64-encoded PDF bytes.

    Raises:
        SandboxExecutionError: If sandbox execution or package install fails.
    """
    md_text = render_markdown(synthesis, critic_output, is_force_delivered)
    md_repr = repr(md_text)
    code = f"""
import re

md = {md_repr}

def md_to_html(text):
    text = re.sub(r'^# (.+)$', r'<h1>\\1</h1>', text, flags=re.MULTILINE)
    text = re.sub(r'^## (.+)$', r'<h2>\\1</h2>', text, flags=re.MULTILINE)
    text = re.sub(r'^### (.+)$', r'<h3>\\1</h3>', text, flags=re.MULTILINE)
    text = re.sub(r'^- (.+)$', r'<li>\\1</li>', text, flags=re.MULTILINE)
    text = re.sub(r'\\*\\*(.+?)\\*\\*', r'<strong>\\1</strong>', text)
    paragraphs = []
    for block in text.split('\\n\\n'):
        block = block.strip()
        if block:
            if block.startswith('<h') or block.startswith('<li'):
                if '<li>' in block:
                    block = '<ul>' + block + '</ul>'
                paragraphs.append(block)
            else:
                paragraphs.append(f'<p>{{block}}</p>')
    return '\\n'.join(paragraphs)

html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
  body {{ font-family: Georgia, serif; max-width: 800px; margin: 40px auto; font-size: 12pt; }}
  h1 {{ font-size: 20pt; }} h2 {{ font-size: 16pt; }} h3 {{ font-size: 13pt; }}
  p, li {{ line-height: 1.6; }}
</style></head>
<body>{{md_to_html(md)}}</body></html>'''

from weasyprint import HTML
HTML(string=html).write_pdf("/tmp/out")
"""
    pdf_bytes = await execute_in_sandbox(code=code, packages=["weasyprint"])
    return _encode_binary_output(pdf_bytes, "pdf")


async def _render_pptx(
    synthesis: SynthesisOutput,
    critic_output: CriticOutput | None,
    is_force_delivered: bool,
) -> str:
    """Render SynthesisOutput as PPTX via sandboxed python-pptx execution.

    Slide layout: title → executive summary → one slide per finding →
    risks → known gaps → references.

    Args:
        synthesis: SynthesisOutput from the coordinator node.
        critic_output: CriticOutput from the critic node, or None.
        is_force_delivered: True when max revisions exhausted without passing.

    Returns:
        Data URI string containing base64-encoded PPTX bytes.

    Raises:
        SandboxExecutionError: If sandbox execution or package install fails.
    """
    doc_data_repr = _build_doc_data(synthesis, critic_output, is_force_delivered)
    code = f"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt

data = json.loads({doc_data_repr})

prs = Presentation()
blank_layout = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]
content_layout = prs.slide_layouts[1]

def add_text_slide(prs, title, body):
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = body
    tf.word_wrap = True

# Title slide
title_slide = prs.slides.add_slide(title_layout)
title_slide.shapes.title.text = data["topic"]
if data["is_force_delivered"] and data["quality_score"] is not None:
    title_slide.placeholders[1].text = (
        f"Quality Notice: Score {{data['quality_score']:.2f}}/5.0"
    )

# Executive summary
add_text_slide(prs, "Executive Summary", data["executive_summary"])

# Key findings — one slide per finding
for f in data["findings"]:
    add_text_slide(prs, f["heading"], f["body"])

# Risks
if data["risks"]:
    add_text_slide(prs, "Risks", "\\n".join(f"• {{r}}" for r in data["risks"]))

# Known gaps
if data["gaps"]:
    add_text_slide(prs, "Known Gaps", "\\n".join(f"• {{g}}" for g in data["gaps"]))

# References
if data["citations"]:
    refs = "\\n".join(
        f"[{{i+1}}] {{c['title']}} — {{c['url']}}"
        for i, c in enumerate(data["citations"])
    )
    add_text_slide(prs, "References", refs)

prs.save("/tmp/out")
"""
    pptx_bytes = await execute_in_sandbox(code=code, packages=["python-pptx"])
    return _encode_binary_output(pptx_bytes, "pptx")


# ---------------------------------------------------------------------------
# Node function
# ---------------------------------------------------------------------------


async def delivery_node(
    state: ResearchAgentState,
) -> dict[str, str]:
    """LangGraph node: format synthesis into the requested output format.

    Binary formats (docx, pdf, pptx) are rendered in an isolated sandbox
    and returned as base64 data URIs stored in final_output.

    Reads state["synthesis"], state["critic_output"], state["output_format"],
    and state["revision_count"].

    Detects force-delivery case (revision_count >= 2 AND critic_output.passed
    == False) and injects a quality caveat at the top of the output.

    Args:
        state: ResearchAgentState. No runtime parameter — no LLM client needed.

    Returns:
        {"final_output": str} — partial state update with rendered document.

    Raises:
        ValueError: If synthesis is None (coordinator must run first).
        SandboxExecutionError: If sandbox rendering fails for binary formats.
    """
    synthesis: SynthesisOutput | None = state["synthesis"]
    if synthesis is None:
        raise ValueError(
            "delivery_node: state['synthesis'] is None — "
            "coordinator must run before delivery."
        )

    critic_output: CriticOutput | None = state["critic_output"]
    output_format = state["output_format"]
    force_delivered = _is_force_delivered(state)

    delivery_start(str(output_format))
    if output_format == "markdown":
        final_output = render_markdown(synthesis, critic_output, force_delivered)
    elif output_format == "docx":
        final_output = await _render_docx(synthesis, critic_output, force_delivered)
    elif output_format == "pdf":
        final_output = await _render_pdf(synthesis, critic_output, force_delivered)
    elif output_format == "pptx":
        final_output = await _render_pptx(synthesis, critic_output, force_delivered)
    else:
        raise NotImplementedError(f"Unknown output_format: {output_format!r}")

    delivery_done(str(output_format), len(final_output))
    return {"final_output": final_output}


__all__ = ["delivery_node", "render_markdown"]
